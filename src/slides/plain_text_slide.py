from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

from typing import List, Callable, Optional
from PIL import Image
import random
import tqdm
import os



from src.font import Font
from .slide_utils import set_shape_transparency, add_paragraph

def generate_plain_text_slide(
    presentation: Presentation,
    title: str,
    text: str,
    font:Font,
    background_path: str = None,
    text_font_coeff:float=0.6,
) -> None:
    """
    Add a slide with title, text placeholders on the blurred background image.

    Args:
        presentation (Presentation): PowerPoint presentation object
        title (str): Title for the slide
        text (str): Text content for the slide
        background_path (str): Path to the background image for the slide
        font (Font): Font object to manage font styles and paths.
        text_font_coeff (float): Coefficient to adjust the font size of the text relative to the title (default is 0.6).
    Returns:
        None
    """

    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)

    slide_height = 9
    slide_width = 16
    margin = min(slide_height, slide_width) / 18

    # Background image
    if background_path:
        pic = slide.shapes.add_picture(
            background_path, 0, 0,
            width=presentation.slide_width,
            height=presentation.slide_height
        )
        # This moves it to the background
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

    # Title placeholder
    title_left = margin
    title_top = margin
    title_width = slide_width - 2 * margin
    title_height = slide_height / 6

    title_box = slide.shapes.add_textbox(
        left=Inches(title_left),
        top=Inches(title_top),
        width=Inches(title_width),
        height=Inches(title_height),
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    
    title_frame.word_wrap = False
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    title_paragraph = add_paragraph(title_frame)
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.text = title

    for max_size in range(font.max_size)[::-5]: 
        try: 
            title_frame.fit_text(
                font_file=font.bold, 
                max_size=max_size, 
                bold=True,
            )
            break
        except: 
            pass
            
    # settings white color and transparency to title shape
    title_fill = title_box.fill
    title_fill.solid()
    title_fill.fore_color.rgb = RGBColor(255, 255, 255)
    set_shape_transparency(title_box, 0.5)

    # Text placeholder
    text_left = Inches(margin)
    text_top = Inches(title_height + margin * 2)
    text_width = Inches(slide_width - 2 * margin)
    text_height = Inches(slide_height - title_height - 3 * margin)
    text_box = slide.shapes.add_textbox(
        left=text_left,
        top=text_top,
        width=text_width,
        height=text_height
    )
    text_frame = text_box.text_frame
    text_frame.clear()
    
    text_frame.word_wrap = False
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    text_paragraph = add_paragraph(text_frame)
    text_paragraph.text = text 
    text_paragraph.alignment = PP_ALIGN.CENTER

    for max_size in range(int(max_size*text_font_coeff))[::-5]: 
        try: 
            text_frame.fit_text(
                font_file=font.basic, 
                max_size=max_size
            )
            break
        except: 
            pass
            
    # Setting text box fill to white with 80% transparency
    text_fill = text_box.fill
    text_fill.solid()
    text_fill.fore_color.rgb = RGBColor(255, 255, 255)
    set_shape_transparency(text_box, 0.5)
    